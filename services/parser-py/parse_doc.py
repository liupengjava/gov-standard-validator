#!/usr/bin/env python3
"""多格式文档解析（PPTX / DOCX / PDF / PPT / DOC）→ 逐页渲染图 + 文本。

统一管线：非 PDF 先用 LibreOffice 转 PDF，再用 PyMuPDF 逐页渲染 PNG + 抽取文本。
PPTX 额外用 python-pptx 补充原生文本（更准）；其余用 fitz 文本。

用法:
  python3 parse_doc.py --in /path/a.pdf --outdir /path/derived/<id> [--dpi 130]
输出: stdout 打印结果 JSON 路径；结果写 outdir/parse_result.json
"""
import argparse, json, os, re, shutil, subprocess, sys, tempfile, uuid, zipfile


def _soffice_to_pdf(path: str, outdir: str) -> str:
    """单次 LibreOffice 转 PDF。soffice 崩溃(嵌入字体等)会以非零退出，由调用方决定是否回退。"""
    profile = f"file://{tempfile.gettempdir()}/lo_profile_{uuid.uuid4().hex}"
    subprocess.run(
        ["soffice", "--headless", "--norestore", f"-env:UserInstallation={profile}",
         "--convert-to", "pdf", "--outdir", outdir, path],
        check=True, capture_output=True, timeout=300,
    )
    pdf = os.path.join(outdir, os.path.splitext(os.path.basename(path))[0] + ".pdf")
    if not os.path.exists(pdf):
        raise FileNotFoundError(f"pdf not produced: {pdf}")
    return pdf


def _strip_embedded_fonts(src: str) -> str:
    """复制 PPTX 并剔除内嵌字体（ppt/fonts/* + presentation.xml 的 embeddedFontLst）。
    某些 deck 的内嵌字体会让 LibreOffice 在 EOT 解码处崩溃(lzcomp 断言)；剔除后回退渲染。
    返回临时 pptx 路径（保持原 basename，使输出 pdf 名一致）。失败抛异常。"""
    tmpdir = tempfile.mkdtemp(prefix="sp_nofont_")
    dst = os.path.join(tmpdir, os.path.basename(src))
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in zin.namelist():
            data = zin.read(n)
            if n.startswith("ppt/fonts/"):
                continue
            if n == "ppt/presentation.xml":
                t = data.decode("utf-8", "ignore")
                t = re.sub(r"<p:embeddedFontLst>.*?</p:embeddedFontLst>", "", t, flags=re.S)
                t = re.sub(r'\sembedTrueTypeFonts="[^"]*"', "", t)
                data = t.encode("utf-8")
            elif n == "ppt/_rels/presentation.xml.rels":
                t = data.decode("utf-8", "ignore")
                t = re.sub(r"<Relationship[^>]*ppt/fonts/[^>]*/>", "", t)
                data = t.encode("utf-8")
            zout.writestr(n, data)
    return dst


def _has_embedded_fonts(path: str) -> bool:
    try:
        with zipfile.ZipFile(path) as z:
            return any(n.startswith("ppt/fonts/") for n in z.namelist())
    except Exception:
        return False


def to_pdf(path: str, outdir: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return path
    # PPTX 内嵌字体会让 LibreOffice 在 EOT 解码处崩溃/卡死(lzcomp 断言 → Abort/ReportCrash 挂起，
    # 进而堆积僵尸 soffice 耗尽内存)。含内嵌字体的 deck 一律先剔除字体再渲染，绕开崩溃路径。
    if ext in (".pptx", ".ppt") and _has_embedded_fonts(path):
        tmp = None
        try:
            tmp = _strip_embedded_fonts(path)
            pdf = _soffice_to_pdf(tmp, outdir)
            sys.stderr.write(f"[to_pdf] 已剔除内嵌字体后渲染(规避 soffice 崩溃): {os.path.basename(path)}\n")
            return pdf
        finally:
            if tmp:
                shutil.rmtree(os.path.dirname(tmp), ignore_errors=True)
    try:
        return _soffice_to_pdf(path, outdir)
    except Exception as e1:
        # 兜底：未预判到的失败，对 PPTX 再试一次剔除字体回退
        if ext in (".pptx", ".ppt"):
            tmp = None
            try:
                tmp = _strip_embedded_fonts(path)
                pdf = _soffice_to_pdf(tmp, outdir)
                sys.stderr.write(f"[to_pdf] 失败回退·剔除内嵌字体渲染成功: {os.path.basename(path)}\n")
                return pdf
            except Exception as e2:
                raise RuntimeError(f"soffice 转 PDF 失败（含剔除字体回退）: {e1} / 回退: {e2}")
            finally:
                if tmp:
                    shutil.rmtree(os.path.dirname(tmp), ignore_errors=True)
        raise


def pptx_extract(path: str):
    """逐页抽取 PPTX 原生：{text, tables(行×列), notes}。失败返回 None。"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        out = []
        for s in prs.slides:
            texts, tables = [], []
            for sh in s.shapes:
                try:
                    if sh.has_text_frame and sh.text_frame.text.strip():
                        texts.append(sh.text_frame.text.strip())
                except Exception:
                    pass
                try:
                    if sh.has_table:
                        rows = [[c.text.strip() for c in r.cells] for r in sh.table.rows]
                        if rows:
                            tables.append(rows)
                except Exception:
                    pass
            notes = ""
            try:
                if s.has_notes_slide:
                    notes = s.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass
            out.append({"text": "\n".join(texts), "tables": tables, "notes": notes})
        return out
    except Exception:
        return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="inp", required=True)
    ap.add_argument("--outdir", required=True)
    ap.add_argument("--dpi", type=int, default=130)
    args = ap.parse_args()
    os.makedirs(args.outdir, exist_ok=True)
    result_path = os.path.join(args.outdir, "parse_result.json")

    import fitz
    try:
        fitz.TOOLS.mupdf_display_errors(False)
    except Exception:
        pass

    pdf = to_pdf(args.inp, args.outdir)
    doc = fitz.open(pdf)
    native = pptx_extract(args.inp) if args.inp.lower().endswith((".pptx", ".ppt")) else None

    slides = []
    for i, page in enumerate(doc):
        png = os.path.join(args.outdir, f"slide_{i+1}.png")
        page.get_pixmap(dpi=args.dpi).save(png)
        text, tables, notes = "", [], ""
        if native and i < len(native):
            text = native[i]["text"]
            tables = native[i]["tables"]
            notes = native[i]["notes"]
        if not text.strip():
            text = page.get_text().strip()
        slides.append({"slide_no": i + 1, "image": png, "text": text, "tables": tables, "notes": notes})
    doc.close()

    json.dump({"ok": True, "pages": len(slides), "slides": slides},
              open(result_path, "w", encoding="utf-8"), ensure_ascii=False)
    print(result_path)


if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(json.dumps({"ok": False, "error": str(e)}, ensure_ascii=False))
        sys.exit(1)
