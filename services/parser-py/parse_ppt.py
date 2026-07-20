#!/usr/bin/env python3
"""把一份 PPTX 解析为逐页结构：渲染图(PNG) + 原生文本/表格/备注。

用法:
  python3 parse_ppt.py --pptx /path/a.pptx --outdir /path/derived/<versionId> [--dpi 130]

输出(stdout, JSON):
  {"ok": true, "pages": N, "slides": [
      {"slide_no":1, "image":"/abs/slide_1.png", "text":"...", "tables":[...], "notes":"..."}, ...]}

说明:
  - 渲染走 LibreOffice(pptx->pdf) + PyMuPDF(pdf->png)，与 python-pptx 的页序按索引对齐。
  - 用独立的 LibreOffice 用户配置目录，允许并发调用。
"""
import argparse, json, os, subprocess, sys, tempfile, uuid


def render_pdf(pptx_path: str, outdir: str) -> str:
    profile = f"file://{tempfile.gettempdir()}/lo_profile_{uuid.uuid4().hex}"
    cmd = [
        "soffice", "--headless", "--norestore",
        f"-env:UserInstallation={profile}",
        "--convert-to", "pdf", "--outdir", outdir, pptx_path,
    ]
    subprocess.run(cmd, check=True, capture_output=True, timeout=300)
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf = os.path.join(outdir, base + ".pdf")
    if not os.path.exists(pdf):
        raise FileNotFoundError(f"pdf not produced: {pdf}")
    return pdf


def render_pngs(pdf_path: str, outdir: str, dpi: int):
    import fitz
    # PyMuPDF 会把 "MuPDF error" 警告打到 stdout，污染 JSON 输出，关闭它
    try:
        fitz.TOOLS.mupdf_display_errors(False)
    except Exception:
        pass
    doc = fitz.open(pdf_path)
    images = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=dpi)
        path = os.path.join(outdir, f"slide_{i+1}.png")
        pix.save(path)
        images.append(path)
    doc.close()
    return images


def extract_pptx(pptx_path: str):
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation(pptx_path)
    slides = []
    for s in prs.slides:
        texts, tables = [], []
        for sh in s.shapes:
            try:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    texts.append(sh.text_frame.text.strip())
            except Exception:
                pass
            try:
                if sh.has_table:
                    rows = []
                    for r in sh.table.rows:
                        rows.append([c.text.strip() for c in r.cells])
                    tables.append(rows)
            except Exception:
                pass
        notes = ""
        try:
            if s.has_notes_slide:
                notes = s.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        slides.append({"text": "\n".join(texts), "tables": tables, "notes": notes})
    return slides


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--outdir", required=True)
    ap.add_argument("--dpi", type=int, default=130)
    ap.add_argument("--result", default=None, help="结果 JSON 写入路径，默认 outdir/parse_result.json")
    args = ap.parse_args()
    os.makedirs(args.outdir, exist_ok=True)
    result_path = args.result or os.path.join(args.outdir, "parse_result.json")

    pdf = render_pdf(args.pptx, args.outdir)
    images = render_pngs(pdf, args.outdir, args.dpi)
    extracted = extract_pptx(args.pptx)

    n = min(len(images), len(extracted))
    slides = []
    for i in range(n):
        slides.append({
            "slide_no": i + 1,
            "image": images[i],
            "text": extracted[i]["text"],
            "tables": extracted[i]["tables"],
            "notes": extracted[i]["notes"],
        })
    result = {"ok": True, "pages": n,
              "pdf_pages": len(images), "pptx_slides": len(extracted),
              "slides": slides}
    with open(result_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False)
    # stdout 只打印结果文件路径，保持干净
    print(result_path)


if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(json.dumps({"ok": False, "error": str(e)}, ensure_ascii=False))
        sys.exit(1)
